"""Generate the Kanad S.H.I.E.L.D. final PPT (11 slides, navy/gold brand).

Run with the system Python (needs python-pptx, qrcode, Pillow — NOT backend deps):
    python backend/scripts/gen_ppt.py
Output: docs/proposals/KanadSHIELD-Final-PPT.pptx
Screenshots + QR codes are read from docs/assets/ppt/ (see scratchpad capture scripts).
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
ASSETS = ROOT / "docs" / "assets" / "ppt"
OUT = ROOT / "docs" / "proposals" / "KanadSHIELD-Final-PPT.pptx"

# ---- brand ----
NAVY = RGBColor(0x0A, 0x11, 0x24)
PANEL = RGBColor(0x11, 0x1B, 0x36)
PANEL_HI = RGBColor(0x18, 0x24, 0x49)
EDGE = RGBColor(0x2A, 0x3A, 0x63)
GOLD = RGBColor(0xF4, 0xB2, 0x3C)
WHITE = RGBColor(0xF5, 0xF7, 0xFC)
SLATE = RGBColor(0x98, 0xA6, 0xC0)
MUTED = RGBColor(0x64, 0x74, 0x94)
GREEN = RGBColor(0x41, 0xD3, 0x92)
RED = RGBColor(0xF8, 0x71, 0x71)
BLUE = RGBColor(0x63, 0xA5, 0xF8)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
SAFFRON = RGBColor(0xFF, 0x99, 0x33)
IND_GREEN = RGBColor(0x13, 0x88, 0x08)

SERIF = "Georgia"
SANS = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)

REPO_URL = "github.com/iampreetdave-max/Guardian-Angle"
VIDEO_URL = "youtu.be/LE9iE1_mCrU"
FOOT = "KANAD S.H.I.E.L.D. 2026  ·  Team VisionScan  ·  Preet Dave"

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide_new():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = NAVY
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=PANEL, edge=EDGE, round_=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = 0.055
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if edge is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = edge
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def run(p, text, size=12, color=WHITE, bold=False, font=SANS, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    r.font.italic = italic
    return r


def para(tf, first=False, space=4, align=PP_ALIGN.LEFT, level=0):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space)
    p.alignment = align
    p.level = level
    return p


def header(s, num, kicker, title, sub=None, tsize=27):
    tf = tb(s, 0.55, 0.30, 12.2, 0.3)
    p = para(tf, first=True)
    run(p, f"{num}  ·  {kicker}", 12.5, GOLD, bold=True)
    tf2 = tb(s, 0.55, 0.58, 12.2, 0.75)
    p2 = para(tf2, first=True)
    run(p2, title, tsize, WHITE, bold=True, font=SERIF)
    if sub:
        tf3 = tb(s, 0.55, 1.22, 12.2, 0.32)
        p3 = para(tf3, first=True)
        run(p3, sub, 12.5, SLATE)


def footer(s, n):
    tf = tb(s, 0.55, 7.08, 9.0, 0.25)
    p = para(tf, first=True)
    run(p, FOOT, 9, MUTED)
    tf2 = tb(s, 12.0, 7.08, 0.8, 0.25)
    p2 = para(tf2, first=True)
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, f"{n} / 11", 9, MUTED)
    # tricolor baseline
    third = 13.333 / 3
    for i, c in enumerate((SAFFRON, WHITE, IND_GREEN)):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(i * third),
                               Inches(7.44), Inches(third), Inches(0.06))
        r.fill.solid()
        r.fill.fore_color.rgb = c
        r.line.fill.background()
        r.shadow.inherit = False


def card(s, x, y, w, h, title=None, lines=None, accent=GOLD, tsize=14,
         bsize=12, pad=0.18, fill=PANEL, lead_space=5):
    box(s, x, y, w, h, fill=fill)
    tf = tb(s, x + pad, y + pad * 0.8, w - 2 * pad, h - 1.4 * pad)
    firstp = True
    if title:
        p = para(tf, first=True, space=lead_space)
        run(p, title, tsize, accent, bold=True)
        firstp = False
    for ln in (lines or []):
        p = para(tf, first=firstp, space=4)
        firstp = False
        if isinstance(ln, tuple):
            label, rest = ln
            run(p, "▸ ", bsize, accent, bold=True)
            run(p, label, bsize, WHITE, bold=True)
            if rest:
                run(p, " — " + rest, bsize, SLATE)
        else:
            run(p, ln, bsize, SLATE)
    return tf


def chip(s, x, y, w, text, color=GOLD, size=10.5, h=0.34, fill=None,
         align=PP_ALIGN.CENTER, bold=True):
    shp = box(s, x, y, w, h, fill=fill if fill else PANEL_HI, edge=color)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run(p, text, size, color, bold=bold)
    return shp


def pic(s, name, x, y, w, caption=None):
    """16:9 screenshot with a hairline border."""
    h = w * 9 / 16
    s.shapes.add_picture(str(ASSETS / name), Inches(x), Inches(y), Inches(w), Inches(h))
    fr = box(s, x, y, w, h, fill=None, edge=EDGE, round_=False)
    fr.line.width = Pt(1.2)
    if caption:
        tf = tb(s, x, y + h + 0.05, w, 0.24)
        p = para(tf, first=True)
        run(p, caption, 10, MUTED)
    return h


def stat(s, x, y, w, h, value, label, note=None, vcolor=GOLD):
    box(s, x, y, w, h)
    tf = tb(s, x + 0.14, y + 0.09, w - 0.28, h - 0.18)
    p = para(tf, first=True, space=1)
    run(p, value, 22, vcolor, bold=True)
    p2 = para(tf, space=0)
    run(p2, label, 10, WHITE, bold=True)
    if note:
        p3 = para(tf, space=0)
        run(p3, note, 9, MUTED)


def arrow(s, x, y, w=0.42, h=0.30, color=GOLD):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                           Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# ================================================================ SLIDE 1
s = slide_new()
chip(s, 0.7, 0.55, 5.05, "KANAD  S.H.I.E.L.D.  2026   ·   CYBER CRIME BRANCH, AHMEDABAD",
     color=GOLD, size=10)
tf = tb(s, 0.7, 1.25, 6.6, 2.1)
p = para(tf, first=True, space=6)
run(p, "VisionScan", 46, WHITE, bold=True, font=SERIF)
run(p, " · CityShield", 46, GOLD, bold=True, font=SERIF)
p = para(tf, space=0)
run(p, "One AI policing platform — from a CCTV frame to a court-ready document.",
    15.5, SLATE)

tf = tb(s, 0.7, 3.55, 6.3, 2.0)
rows = [
    ("Team", "Team VisionScan"),
    ("Team member", "Preet Dave  (solo build)"),
    ("Organization", "Independent final-year student project"),
    ("Built for", "Cyber Crime Branch, Ahmedabad City Police"),
]
first = True
for k, v in rows:
    p = para(tf, first=first, space=6)
    first = False
    run(p, f"{k:<14}", 12.5, GOLD, bold=True)
    run(p, "   " + v, 13, WHITE)
p = para(tf, space=0)
run(p, "CODE  ", 12, GREEN, bold=True)
run(p, "https://" + REPO_URL, 12, GOLD, bold=True)

s.shapes.add_picture(str(ASSETS / "qr_github.png"), Inches(0.7), Inches(5.72),
                     Inches(1.0), Inches(1.0))
tf = tb(s, 1.85, 6.02, 4.9, 0.8)
p = para(tf, first=True, space=2)
run(p, "Scan for the full source + demo video", 11, SLATE)
p = para(tf, space=0)
run(p, "Runs offline on this laptop; a hosted instance can be started on request.",
    9.5, MUTED)

# right collage
pic(s, "04_citymap_risk.png", 7.35, 1.30, 5.35)
pic(s, "07_livealerts.png", 6.85, 4.05, 3.6)
pic(s, "11_crimegpt.png", 10.6, 4.42, 2.35)
tf = tb(s, 7.35, 6.30, 5.4, 0.3)
p = para(tf, first=True)
run(p, "Running on a field laptop today — predictive City Map, Live Alerts, CrimeGPT", 10.5, MUTED)
footer(s, 1)

# ================================================================ SLIDE 2
s = slide_new()
header(s, "02", "PROBLEM STATEMENT",
       "Investigators drown in footage, paperwork and siloed tools",
       "Flagship statement: GIS-enabled Crime Hotspot Mapping & Predictive Patrol Routing — Cyber Crime Branch, Ahmedabad (Category II).",
       tsize=24)

card(s, 0.55, 1.75, 6.05, 2.6, "The problem", [
    ("Manual CCTV review", "one incident can mean hours of scrubbing footage across cameras — evidence is missed or found late."),
    ("Gut-feel patrolling", "beat allocation rarely uses data; hotspots shift faster than monthly review meetings."),
    ("Paper-heavy casework", "FIRs, panchnamas, custody letters are drafted by hand under BNSS time limits."),
    ("Cyber-fraud golden hour", "1930 escalation is most effective in the first minutes — intake is slow and unstructured."),
])
card(s, 6.85, 1.75, 5.9, 2.6, "Why this matters", [
    ("Response time saves outcomes", "faster detection and dispatch directly cut harm and raise recovery odds."),
    ("Force multiplication", "Ahmedabad-scale cities cannot staff every screen; AI watches feeds continuously."),
    ("Trust & accountability", "predictive policing without transparency invites bias and legal challenge."),
    ("Five statements, one root cause", "detection, prediction, documentation and intel all fail apart when data is siloed."),
])

tf = tb(s, 0.55, 4.52, 12.2, 0.3)
p = para(tf, first=True)
run(p, "Existing challenges", 13.5, GOLD, bold=True)
run(p, "     Aligned with Gujarat Police Innovation Challenge 2026 (80,000+ CCTVs, 17 Aug) "
       "and e-Zero FIR with I4C (27 Jul)", 10, SLATE)
ch = [
    ("Proprietary & costly", "PredPol / Briefcam-class tools are licensed, closed and priced for large forces."),
    ("Cloud & GPU dependent", "footage leaves the premises; connectivity and hardware become blockers."),
    ("Black-box models", "opaque scoring (the COMPAS criticism) — no 'why' behind a hotspot."),
    ("Fragmented products", "detection, prediction, dispatch and drafting live in 4–5 separate databases."),
]
for i, (t, b) in enumerate(ch):
    card(s, 0.55 + i * 3.11, 4.88, 2.96, 1.55, t, [b], tsize=12, bsize=10.5)
footer(s, 2)

# ================================================================ SLIDE 3
s = slide_new()
header(s, "03", "PROPOSED SOLUTION", "One platform. Five modules. One closed loop.",
       "A single codebase, database and security layer answers five Category-II problem statements at once.")

mods = [
    ("VisionScan", "Ctrl-F for CCTV", "Natural-language, object, image and suspect-face search over indexed footage (CLIP + YOLOv8 + ArcFace)."),
    ("Anomaly Watch", "Always-on detection", "Every feed scanned for fire, smoke, accidents, weapons, violence — debounced alerts with keyframes."),
    ("City Map", "Predict & patrol", "Recency-weighted hotspot forecast over 30 Ahmedabad localities + optimized patrol routes (NN + 2-opt)."),
    ("CrimeGPT", "Facts once → every document", "Suggests BNS/BNSS sections; generates 7 court-ready PDFs in English, Hindi, Gujarati."),
    ("Arbiter + GovIntel", "Legal intelligence", "Offline section lookup, FIR drafts, legal Q&A + unified GR/Act/judgment search and 1930 cyber intake."),
]
pos = [(0.55, 1.85), (4.66, 1.85), (8.77, 1.85), (2.6, 3.85), (6.71, 3.85)]
for (name, tag, body), (x, y) in zip(mods, pos):
    box(s, x, y, 4.0, 1.75)
    tf = tb(s, x + 0.18, y + 0.14, 3.64, 1.55)
    p = para(tf, first=True, space=1)
    run(p, name, 14.5, WHITE, bold=True)
    p = para(tf, space=4)
    run(p, tag.upper(), 9.5, GOLD, bold=True)
    p = para(tf, space=0)
    run(p, body, 11, SLATE)

vals = ["Closed loop: detect → case → risk → dispatch",
        "Offline-capable · CPU-only · field-ready",
        "Zero paid APIs · ₹0 infrastructure",
        "One database · RBAC · audit logs"]
for i, v in enumerate(vals):
    chip(s, 0.55 + i * 3.11, 6.25, 2.96, v, color=GOLD, size=9, h=0.5)
tf = tb(s, 0.55, 6.02, 12.2, 0.25)
p = para(tf, first=True)
run(p, "Value proposition", 12, GOLD, bold=True)
footer(s, 3)

# ================================================================ SLIDE 4
s = slide_new()
header(s, "04", "SYSTEM ARCHITECTURE", "From camera frame to dispatched unit",
       "High-level architecture and data flow — every module reads and writes one shared evidence store.")

cols = [
    ("SOURCES", BLUE, [
        "CCTV clips & live feeds",
        "Citizen complaints portal",
        "NCRP cyber-fraud intake (1930)",
        "Gov notifications / RSS",
    ]),
    ("AI CORE", GOLD, [
        "CLIP ViT-B/32 zero-shot vision",
        "YOLOv8 objects + fire/smoke",
        "ArcFace watchlist matching",
        "Risk model + 2-opt routing",
        "MiniLM + ChromaDB legal RAG",
    ]),
    ("PLATFORM", GREEN, [
        "FastAPI + SQLite evidence store",
        "JWT · RBAC · OWASP middleware",
        "Cases, dispatch & notifications",
        "Audit log · SHA-256 exports",
    ]),
    ("INTERFACES", VIOLET, [
        "React command dashboard",
        "GIS City Map (Leaflet/OSM)",
        "Citizen & field-mobile views",
        "Branded PDF documents",
    ]),
]
cw, cy, chh = 2.82, 1.8, 3.0
for i, (t, c, items) in enumerate(cols):
    x = 0.55 + i * (cw + 0.42)
    box(s, x, cy, cw, chh)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(cy),
                             Inches(cw), Inches(0.09))
    bar.fill.solid(); bar.fill.fore_color.rgb = c
    bar.line.fill.background(); bar.shadow.inherit = False
    tf = tb(s, x + 0.16, cy + 0.2, cw - 0.32, chh - 0.35)
    p = para(tf, first=True, space=6)
    run(p, t, 13, c, bold=True)
    for it in items:
        p = para(tf, space=3.5)
        run(p, "▸ ", 11, c, bold=True)
        run(p, it, 11, WHITE)
    if i < 3:
        arrow(s, x + cw + 0.045, cy + chh / 2 - 0.15)

# closed loop
tf = tb(s, 0.55, 5.05, 12.2, 0.28)
p = para(tf, first=True)
run(p, "The closed loop (live workflow)", 12.5, GOLD, bold=True)
loop = ["Anomaly detected on a feed", "Case auto-created + keyframe evidence",
        "Locality risk score bumped", "Nearest patrol unit dispatched"]
lw = 2.78
for i, step in enumerate(loop):
    x = 0.55 + i * (lw + 0.37)
    shp = box(s, x, 5.42, lw, 0.78, fill=PANEL_HI, edge=GOLD)
    t2 = shp.text_frame
    t2.word_wrap = True
    t2.vertical_anchor = MSO_ANCHOR.MIDDLE
    t2.margin_left = t2.margin_right = Inches(0.1)
    pp = t2.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    run(pp, f"{i+1}.  ", 11, GOLD, bold=True)
    run(pp, step, 11, WHITE, bold=True)
    if i < 3:
        arrow(s, x + lw + 0.02, 5.68, w=0.33, h=0.26)
tf = tb(s, 0.55, 6.38, 12.2, 0.3)
p = para(tf, first=True)
run(p, "Data flow: sources → AI core → one SQLite store → dashboards & documents. "
       "Optional pieces (Gemini, map tiles, gov RSS) degrade gracefully — the loop above runs fully offline.",
    10.5, SLATE)
footer(s, 4)

# ================================================================ SLIDE 5
s = slide_new()
header(s, "05", "TECHNOLOGY STACK", "Open models, boring infrastructure, zero paid APIs",
       "Everything below is open-source or free-tier — the platform runs on hardware a police station already owns.")

stacks = [
    ("Software", [
        ("Backend", "Python · FastAPI · SQLite · Uvicorn"),
        ("Frontend", "React · Vite · Tailwind CSS"),
        ("GIS", "Leaflet + OpenStreetMap · Recharts"),
        ("Packaging", "Docker Compose · one-command start"),
        ("Testing", "pytest + httpx — 81-test suite"),
    ]),
    ("AI / ML models", [
        ("Vision", "OpenAI CLIP ViT-B/32 (zero-shot)"),
        ("Detection", "Ultralytics YOLOv8 + fire/smoke model"),
        ("Faces", "ArcFace / InsightFace watchlists"),
        ("Legal RAG", "MiniLM embeddings + ChromaDB"),
        ("Prediction", "Recency-weighted risk + NN/2-opt routing"),
    ]),
    ("Hardware", [
        ("Runs on", "any CPU-only laptop / desktop"),
        ("Reference host", "2 vCPU · 8 GB RAM VM — no GPU anywhere"),
        ("No GPU needed", "clip indexing ~30–90 s per clip"),
        ("Scale-up path", "single GPU → real-time RTSP feeds"),
        ("Field mode", "fully offline on a station laptop"),
    ]),
    ("Cloud & APIs", [
        ("Hosting", "any commodity VM or a field laptop (₹0 cost)"),
        ("HTTPS", "Caddy auto-TLS · git-push deploys"),
        ("Optional", "Gemini — falls back to offline templates"),
        ("Map tiles", "OSM — cached, degrades gracefully"),
        ("No lock-in", "no paid or proprietary API anywhere"),
    ]),
]
for i, (t, lines) in enumerate(stacks):
    card(s, 0.55 + i * 3.11, 1.8, 2.96, 3.6, t, lines, tsize=13.5, bsize=10.3)

chip(s, 0.55, 5.7, 5.9, "Security: JWT + RBAC · OWASP Top-10 middleware · rate limiting · audit logging",
     color=GREEN, size=10, h=0.45, align=PP_ALIGN.LEFT)
chip(s, 6.65, 5.7, 6.1, "Compliance: MIT-licensed code · third-party notices · AGPL swap path documented",
     color=BLUE, size=10, h=0.45, align=PP_ALIGN.LEFT)
footer(s, 5)

# ================================================================ SLIDE 6
s = slide_new()
header(s, "06", "DEMONSTRATION", "It runs live — try it during judging",
       "Screenshots below were captured from the running platform (synthetic demo data, no real persons or cases).")

pic(s, "07_livealerts.png", 0.55, 1.8, 3.72, caption="Live Alerts — fire/smoke/violence with confidence")
pic(s, "09_visionscan_search.png", 4.42, 1.8, 3.72, caption='VisionScan — "flames engulfing a building" → ranked moments')
pic(s, "05_citymap_routes.png", 0.55, 4.35, 3.72, caption="City Map — hotspots + optimized 2-unit patrol plan")
pic(s, "11_crimegpt.png", 4.42, 4.35, 3.72, caption="CrimeGPT — 7 court-ready PDFs · en / hi / gu")

sx = 8.5
stat(s, sx, 1.80, 2.03, 1.05, "0.790", "Hit-Rate@10", "90% CI 0.77–0.81")
stat(s, sx + 2.13, 1.80, 2.03, 1.05, "2.37×", "PAI@10", "94% of 2.53× oracle")
stat(s, sx, 2.98, 2.03, 1.05, "79%", "next-week crime", "in top-10 zones (⅓ of city)")
stat(s, sx + 2.13, 2.98, 2.03, 1.05, "2 / 2", "planted surges", "caught in live top-10")
tf = tb(s, sx, 4.12, 4.2, 0.5)
p = para(tf, first=True)
run(p, "Rolling-origin backtest vs frequency / prior / random baselines · 81 backend tests passing.",
    9.5, SLATE)

box(s, sx, 4.66, 4.16, 2.06)
tf = tb(s, sx + 0.16, 4.78, 3.9, 1.9)
p = para(tf, first=True, space=3)
run(p, "DEMO ACCESS", 11, GOLD, bold=True)
p = para(tf, space=3)
run(p, "docker compose up -d  ->  localhost:8080", 11, WHITE, bold=True)
creds = [("Admin", "admin@city.gov / admin123"),
         ("Lead", "lead@city.gov / lead123"),
         ("Officer", "officer@city.gov / officer123"),
         ("Citizen", "citizen@example.com / citizen123")]
for role, cred in creds:
    p = para(tf, space=1.5)
    run(p, f"{role:<9}", 9.5, GOLD, bold=True, font="Consolas")
    run(p, cred, 9.5, SLATE, font="Consolas")
footer(s, 6)

# ================================================================ SLIDE 7
s = slide_new()
header(s, "07", "INNOVATION", "What no existing tool does",
       "Novelty is not a bigger model — it is the loop, the honesty and the sovereignty.")

nov = [
    ("The closed loop", "A live detection becomes a geo-tagged case with keyframe evidence, bumps that locality's risk and pings the nearest unit — detection becomes action, in one system, automatically."),
    ("Honest, measured AI", "We publish backtests, confidence intervals and the oracle ceiling (best possible score on the data) — plus a 'Why this hotspot?' breakdown for every prediction. No black box."),
    ("Offline-first sovereignty", "Footage and models never leave the laptop. Every cloud dependency is optional and degrades gracefully — the platform demos with Wi-Fi unplugged."),
]
for i, (t, b) in enumerate(nov):
    card(s, 0.55 + i * 4.11, 1.75, 3.96, 1.72, t, [b], tsize=13, bsize=10.5)

rows = [
    ("Capability", "VisionScan · CityShield", "PredPol-class", "Briefcam-class", "Legal-AI tools"),
    ("Hotspot forecast + patrol routing", "✓", "✓", "—", "—"),
    ("CCTV semantic search + live anomaly alerts", "✓", "—", "✓", "—"),
    ("BNS/BNSS drafting, FIR-ready PDFs", "✓", "—", "—", "✓"),
    ("Works offline, CPU-only, no licence fees", "✓", "—", "—", "—"),
    ("Transparent scoring + published backtests", "✓", "—", "—", "—"),
    ("One shared database across all of the above", "✓", "—", "—", "—"),
]
tbl_shape = s.shapes.add_table(len(rows), 5, Inches(0.55), Inches(3.75),
                               Inches(12.23), Inches(2.9))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(4.63)
for ci in range(1, 5):
    tbl.columns[ci].width = Inches(1.9)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_HI if ri == 0 else PANEL
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfc = cell.text_frame
        tfc.word_wrap = True
        pc = tfc.paragraphs[0]
        pc.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
        if ri == 0:
            run(pc, val, 10.5, GOLD if ci == 1 else WHITE, bold=True)
        elif ci == 0:
            run(pc, val, 10.5, WHITE)
        else:
            good = val == "✓"
            run(pc, val, 12, (GREEN if ci == 1 else SLATE) if good else MUTED,
                bold=good)
footer(s, 7)

# ================================================================ SLIDE 8
s = slide_new()
header(s, "08", "IMPACT", "Fewer blind spots, faster response, less paperwork",
       "Quantified on the live demo build — methodology transfers directly to real district data under an MoU.")

card(s, 0.55, 1.78, 6.05, 2.42, "For law enforcement", [
    ("Patrol smarter", "covering the top-10 predicted zones (⅓ of the city) captured 79% of next-week crime in backtests."),
    ("CCTV in seconds", "semantic search replaces hours of manual footage scrubbing."),
    ("Paperwork in minutes", "7 court-ready documents generated from facts entered once."),
    ("Always-on eyes", "anomaly watch monitors every feed a human cannot."),
])
card(s, 6.85, 1.78, 5.93, 2.42, "For public safety", [
    ("Faster response", "auto-dispatch of the nearest unit the moment an incident is detected."),
    ("Cyber golden hour", "structured 1930 intake pushes victims to act in the recovery window."),
    ("Citizen participation", "complaint portal with live status — transparency builds trust."),
    ("City-wide advisories", "severity-graded disaster broadcasts reach every user."),
])
card(s, 0.55, 4.42, 6.05, 2.3, "Expected outcomes", [
    ("Measurable patrol lift", "2.37× more crime coverage than uniform patrolling (PAI@10)."),
    ("Evidence integrity", "SHA-256-stamped exports + audit trail stand up in court."),
    ("Institutional memory", "every complaint, case, detection and document in one queryable store."),
])
card(s, 6.85, 4.42, 5.93, 2.3, "Use cases", [
    ("Station command", "dashboard + City Map for daily beat planning."),
    ("Control room", "Live Alerts wall for multi-feed monitoring."),
    ("Field officer", "mobile view — cases, dispatch and evidence on the move."),
    ("Investigating officer", "VisionScan search + CrimeGPT drafting on any case."),
])
footer(s, 8)

# ================================================================ SLIDE 9
s = slide_new()
header(s, "09", "SCALABILITY & FUTURE SCOPE", "Pilot-ready today, district-scale by design",
       "The demo is the deployment — the same containers move from a student VM to a station server unchanged.")

card(s, 0.55, 1.78, 6.05, 2.35, "Deployment strategy", [
    ("Today", "one-command Docker Compose deploy (Caddy auto-HTTPS); pilot ran on an Azure student VM until the credit lapsed."),
    ("Pilot", "one on-prem station server — data never leaves the premises (DPDP-friendly)."),
    ("One command", "start.ps1 / docker compose up -d brings up the full platform."),
    ("Real data", "ingest standard incident schemas under an MoU with the district."),
])
card(s, 6.85, 1.78, 5.93, 2.35, "Scalability", [
    ("Storage", "SQLite → PostgreSQL; footage → object storage when volume demands."),
    ("Compute", "one GPU unlocks real-time RTSP camera ingestion."),
    ("Search", "FAISS index → shared vector DB across stations."),
    ("Federation", "multi-station / multi-district rollout on the same RBAC model."),
])
card(s, 0.55, 4.35, 6.05, 2.35, "Future enhancements", [
    ("Sharper detection", "specialized fire/smoke localization already in the pipeline (next deploy)."),
    ("Drift monitoring", "built-in rolling-origin backtests re-run on fresh data to catch model drift."),
    ("Legal library growth", "lawyer-reviewed BNS/BNSS pattern set expands with practitioner feedback."),
    ("ANPR & watchlists", "plate recognition and richer suspect tooling as opt-in modules."),
])
card(s, 6.85, 4.35, 5.93, 2.35, "Sustainability", [
    ("No licence fees", "open models + MIT-licensed code — cost stays near zero at pilot scale."),
    ("Commodity hardware", "runs on machines departments already own; no cloud commitment."),
    ("Human-in-the-loop", "system suggests, officers decide — with audit logs for accountability."),
    ("Open compliance", "third-party notices shipped; AGPL detector swap path documented."),
])
footer(s, 9)

# ================================================================ SLIDE 10
s = slide_new()
header(s, "10", "SUPPORTING DOCUMENTS & RESOURCES", "Everything a judge needs, one scan away",
       "Every QR is a permanent link; all PDFs are bundled with the submission.")

res = [
    ("qr_video.png", "Demo video (2:46)", VIDEO_URL, RED),
    ("qr_github.png", "Source code (GitHub)", REPO_URL, GOLD),
    ("qr_drive.png", "Test clips + demo logins", "Google Drive folder", VIOLET),
]
for i, (qr, title, url, c) in enumerate(res):
    x = 0.55 + i * 4.13
    box(s, x, 1.78, 3.98, 2.75)
    s.shapes.add_picture(str(ASSETS / qr), Inches(x + 1.34), Inches(1.98),
                         Inches(1.3), Inches(1.3))
    tf = tb(s, x + 0.15, 3.42, 3.68, 1.1)
    p = para(tf, first=True, space=2)
    p.alignment = PP_ALIGN.CENTER
    run(p, title, 12, c, bold=True)
    p = para(tf, space=0)
    p.alignment = PP_ALIGN.CENTER
    run(p, url, 9, SLATE)

docs = [
    ("📄  Project report & 5 proposal PDFs", "branded, bundled in the submission ZIP"),
    ("🎥  Demo video", "captioned 2:46 walkthrough — youtu.be/LE9iE1_mCrU (QR above)"),
    ("📊  Dataset documentation", "docs/AHMEDABAD_CRIME_DATA.md — NCRB-derived synthetic data, full provenance"),
    ("✅  Validation & security reports", "docs/VALIDATION.md · security testing report · 81-test suite"),
]
for i, (t, b) in enumerate(docs):
    x = 0.55 + (i % 2) * 6.19
    y = 4.85 + (i // 2) * 0.95
    shp = box(s, x, y, 6.04 if i % 2 == 0 else 6.18, 0.82)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.12)
    p = tf.paragraphs[0]
    run(p, t + "   ", 12, WHITE, bold=True)
    run(p, b, 10.5, SLATE)
footer(s, 10)

# ================================================================ SLIDE 11
s = slide_new()
tf = tb(s, 0.7, 1.7, 7.6, 2.2)
p = para(tf, first=True, space=6)
run(p, "Thank you.", 44, WHITE, bold=True, font=SERIF)
p = para(tf, space=0)
run(p, "Questions & answers — everything shown today runs on this laptop, offline.",
    15, SLATE)

box(s, 0.7, 3.9, 6.6, 2.3)
tf = tb(s, 0.95, 4.12, 6.1, 1.9)
p = para(tf, first=True, space=6)
run(p, "CONTACT", 11, GOLD, bold=True)
p = para(tf, space=3)
run(p, "Preet Dave  ·  Team VisionScan", 15, WHITE, bold=True)
p = para(tf, space=3)
run(p, "iampreetdave@gmail.com", 13, GOLD)
p = para(tf, space=0)
run(p, "Built for the Cyber Crime Branch, Ahmedabad City Police — Kanad S.H.I.E.L.D. 2026",
    11.5, SLATE)

box(s, 8.6, 2.4, 3.9, 3.8)
s.shapes.add_picture(str(ASSETS / "qr_github.png"), Inches(9.55), Inches(2.75),
                     Inches(2.0), Inches(2.0))
tf = tb(s, 8.8, 4.88, 3.5, 1.4)
p = para(tf, first=True, space=2)
p.alignment = PP_ALIGN.CENTER
run(p, "Scan — run it yourself", 12, WHITE, bold=True)
p = para(tf, space=0)
p.alignment = PP_ALIGN.CENTER
run(p, REPO_URL + "\nDemo video: " + VIDEO_URL + "\nadmin@city.gov / admin123", 10.5, SLATE)
footer(s, 11)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides)")
